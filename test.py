import streamlit as st
from pinecone import Pinecone, ServerlessSpec
from youtube_transcript_api import YouTubeTranscriptApi
from PyPDF2 import PdfReader
from pptx import Presentation
from sentence_transformers import SentenceTransformer
from selenium import webdriver
from selenium.webdriver.chrome.service import Service
from selenium.webdriver.chrome.options import Options
from bs4 import BeautifulSoup
import requests

# Pinecone-configuratie
PINECONE_API_KEY = "pcsk_4v8YTF_6Sgtnwh2Tm38koMurffJJLUp4eyHncT843KmkeKN3GVbjqSNbFPtzjBiDTfkF6V"
INDEX_NAME = "projectvito"

@st.cache_resource
def initialize_pinecone(api_key, index_name):
    pc = Pinecone(api_key=api_key)
    return pc.Index(index_name)

index = initialize_pinecone(PINECONE_API_KEY, INDEX_NAME)


# Embedder configuratie
@st.cache_resource
def load_embedder():
    return SentenceTransformer("all-MiniLM-L6-v2")
embedder = load_embedder()


# Functie om PowerPoint-bestanden te verwerken
def extract_text_from_pptx(file):
    """Haal tekst uit een PowerPoint-bestand (.pptx)."""
    try:
        presentation = Presentation(file)
        full_text = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except Exception as e:
        print(f"Fout bij verwerken van PowerPoint: {e}")
        return None

# Functie om webpagina's te scrapen
def scrape_webpage(url):
    try:
        # Probeer eerst met requests
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        soup = BeautifulSoup(response.text, "html.parser")
        text = soup.get_text(strip=True)
        return text
    except Exception as e:
        print(f"Requests-fout: {e}")
        return None

# Functie om tekst te limiteren voor Pinecone metadata
def limit_text_size(text, max_bytes=40000):
    """Limiteer tekst tot een specifiek aantal bytes voor Pinecone metadata."""
    text_bytes = text.encode('utf-8')
    if len(text_bytes) > max_bytes:
        # Verkort de tekst tot we onder de limiet zitten
        while len(text_bytes) > max_bytes:
            text = text[:len(text)-1000]  # Verwijder 1000 karakters per keer
            text_bytes = text.encode('utf-8')
    return text


# Tabs voor upload-opties
tab1, tab2, tab3, tab4 = st.tabs(["📄 Bestanden", "📹 YouTube-video's", "🌐 Webpagina's", "📊 PowerPoints"])

# Tab 1: PDF-bestanden
with tab1:
    st.header("Upload een PDF-bestand")
    uploaded_file = st.file_uploader("Kies een PDF-bestand", type=["pdf"])

    if uploaded_file:
        try:
            pdf_reader = PdfReader(uploaded_file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()

            # Limiteer de tekstgrootte
            limited_text = limit_text_size(text)
            
            # Genereer embedding
            embedding = embedder.encode(limited_text).tolist()

            # Upsert naar Pinecone
            index.upsert([(uploaded_file.name, embedding, {"type": "PDF", "text": limited_text})])
            st.success("PDF succesvol toegevoegd aan Pinecone!")
        except Exception as e:
            st.error(f"Fout bij verwerken van PDF: {e}")

# Tab 2: YouTube-video's
with tab2:
    st.header("Voeg een YouTube-video toe")
    youtube_url = st.text_input("Voer een YouTube-video-URL in:")

    if st.button("YouTube toevoegen"):
        if youtube_url:
            try:
                video_id = youtube_url.split("v=")[1].split("&")[0]
                transcript = YouTubeTranscriptApi.get_transcript(video_id)
                full_text = " ".join([t["text"] for t in transcript])

                # Limiteer de tekstgrootte
                limited_text = limit_text_size(full_text)
                
                # Genereer embedding
                embedding = embedder.encode(limited_text).tolist()

                # Upsert naar Pinecone
                index.upsert([(video_id, embedding, {"type": "YouTube", "text": limited_text})])
                st.success("YouTube-video succesvol toegevoegd aan Pinecone!")
            except Exception as e:
                st.error(f"Fout bij ophalen van transcript: {e}")
        else:
            st.warning("Voer een geldige YouTube-URL in.")

# Tab 3: Webpagina's
with tab3:
    st.header("Voeg een webpagina toe")
    webpage_url = st.text_input("Voer een webpagina-URL in:")

    if st.button("Webpagina toevoegen"):
        if webpage_url:
            try:
                text = scrape_webpage(webpage_url)
                if text:
                    # Limiteer de tekstgrootte
                    limited_text = limit_text_size(text)
                    
                    # Genereer embedding
                    embedding = embedder.encode(limited_text).tolist()

                    # Upsert naar Pinecone
                    index.upsert([(webpage_url, embedding, {"type": "Webpage", "text": limited_text})])
                    st.success("Webpagina succesvol toegevoegd aan Pinecone!")
                else:
                    st.warning("Kon geen tekst van de webpagina ophalen.")
            except Exception as e:
                st.error(f"Fout bij verwerken van webpagina: {e}")
        else:
            st.warning("Voer een geldige webpagina-URL in.")

# Tab 4: PowerPoints
with tab4:
    st.header("Upload een PowerPoint-bestand")
    pptx_file = st.file_uploader("Kies een PowerPoint-bestand", type=["pptx"])

    if pptx_file:
        try:
            text = extract_text_from_pptx(pptx_file)
            if text:
                # Limiteer de tekstgrootte
                limited_text = limit_text_size(text)
                
                # Genereer embedding
                embedding = embedder.encode(limited_text).tolist()

                # Upsert naar Pinecone
                index.upsert([(pptx_file.name, embedding, {"type": "PowerPoint", "text": limited_text})])
                st.success("PowerPoint succesvol toegevoegd aan Pinecone!")
            else:
                st.warning("Geen tekst gevonden in het PowerPoint-bestand.")
        except Exception as e:
            st.error(f"Fout bij verwerken van PowerPoint: {e}")
